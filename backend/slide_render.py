"""SVG + pptx rendering for single-slide image+text layouts.

Reuses TextItem rendering from the diagram renderers; adds image drawing.
Images embed as base64 in SVG (self-contained preview) and add_picture in pptx.
"""
from __future__ import annotations

import base64
import io
import mimetypes
import os

from pptx import Presentation
from pptx.util import Inches

from layout import SLIDE_W, SLIDE_H, TextItem
from slide_layout import ImageItem, layout_slide_layout
from render_svg import PX, W, H, _text_item_svg
from export_pptx import _add_text
from models import SlideLayout

_PLACE_FILL = "#eef1f6"
_PLACE_STROKE = "#c3cad4"


def _data_uri(path: str) -> str | None:
    try:
        mime = mimetypes.guess_type(path)[0] or "image/png"
        with open(path, "rb") as f:
            return f"data:{mime};base64," + base64.b64encode(f.read()).decode()
    except Exception:
        return None


def _image_svg(it: ImageItem) -> str:
    x, y, w, h = it.x * PX, it.y * PX, it.w * PX, it.h * PX
    uri = _data_uri(it.path) if it.path else None
    if uri:
        return (f'<image x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                f'href="{uri}" preserveAspectRatio="xMidYMid meet"/>')
    return (f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
            f'fill="{_PLACE_FILL}" stroke="{_PLACE_STROKE}" stroke-dasharray="6 4"/>'
            f'<text x="{x+w/2:.1f}" y="{y+h/2:.1f}" text-anchor="middle" '
            f'fill="#8a93a0" font-size="16">[이미지 {it.ref}]</text>')


def render_slide_layout_svg(layout: SlideLayout, assets: dict) -> str:
    body = []
    for it in layout_slide_layout(layout, assets):
        if isinstance(it, TextItem):
            body.append(_text_item_svg(it))
        elif isinstance(it, ImageItem):
            body.append(_image_svg(it))
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" width="100%" '
        f'style="background:#fff;border:1px solid #d7dbe3;border-radius:8px;display:block;">'
        f'{"".join(body)}</svg>'
    )


def build_slide_pptx(layout: SlideLayout, assets: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for it in layout_slide_layout(layout, assets):
        if isinstance(it, TextItem):
            _add_text(slide, it)
        elif isinstance(it, ImageItem):
            if it.path and os.path.exists(it.path):
                slide.shapes.add_picture(it.path, Inches(it.x), Inches(it.y),
                                         Inches(it.w), Inches(it.h))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
