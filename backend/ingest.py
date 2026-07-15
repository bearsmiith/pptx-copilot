"""Attachment ingestion for generation requests.

Uploads are saved to a per-request directory; the claude_cli provider lets
the model Read them directly (images via vision, svg/text as text). Office
binaries (pptx/docx/xlsx) are pre-extracted to a .txt sidecar since the CLI
Read tool cannot parse them.
"""
from __future__ import annotations

import os
import re
import uuid

UPLOAD_ROOT = "/tmp/pptx-copilot-uploads"
MAX_FILE_MB = 15
MAX_FILES = 8

_IMAGE = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
_TEXTY = {".svg", ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
          ".html", ".py", ".js", ".ts", ".log", ".tsv"}
_OFFICE = {".pptx", ".docx", ".xlsx"}

IMAGE_MIME = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
              ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"}


def _safe_name(name: str) -> str:
    base = os.path.basename(name or "file")
    base = re.sub(r"[^\w.\-가-힣 ]", "_", base).strip() or "file"
    return base[:80]


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    parts = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        parts.append(f"--- slide {i+1} ---")
        for shp in slide.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def _extract_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- sheet: {ws.title} ---")
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r > 200:
                parts.append("... (truncated)")
                break
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    r = PdfReader(path)
    return "\n".join((p.extract_text() or "") for p in r.pages[:40])


def save_attachments(files: list[tuple[str, bytes]]) -> tuple[str, list[dict]]:
    """Save (filename, data) pairs. Returns (workdir, manifest).

    manifest entries:
      name       original filename (sanitized)
      kind       image | text | office | pdf | skip
      read_path  path for the claude Read tool (None if skip)
      text_path  extracted-text sidecar (office/pdf) for text-only providers
      note       human hint or skip reason
    """
    if len(files) > MAX_FILES:
        raise ValueError(f"최대 {MAX_FILES}개 파일까지 첨부 가능합니다")
    workdir = os.path.join(UPLOAD_ROOT, uuid.uuid4().hex[:12])
    os.makedirs(workdir, exist_ok=True)
    manifest: list[dict] = []

    def entry(**kw):
        base = {"name": None, "kind": "text", "read_path": None,
                "text_path": None, "note": None}
        base.update(kw)
        manifest.append(base)

    for name, data in files:
        if len(data) > MAX_FILE_MB * 1024 * 1024:
            raise ValueError(f"{name}: {MAX_FILE_MB}MB 초과")
        safe = _safe_name(name)
        path = os.path.join(workdir, safe)
        with open(path, "wb") as f:
            f.write(data)
        ext = os.path.splitext(safe)[1].lower()

        if ext in _IMAGE:
            entry(name=safe, kind="image", read_path=path)
        elif ext in _OFFICE or ext == ".pdf":
            extractor = {".pptx": _extract_pptx, ".docx": _extract_docx,
                         ".xlsx": _extract_xlsx, ".pdf": _extract_pdf}[ext]
            try:
                text = extractor(path)
            except Exception as e:
                text = f"(extraction failed: {e})"
            side = path + ".extracted.txt"
            with open(side, "w", encoding="utf-8") as f:
                f.write(text[:200_000])
            if ext == ".pdf":
                # claude reads the pdf natively; text sidecar for others
                entry(name=safe, kind="pdf", read_path=path, text_path=side)
            else:
                entry(name=safe, kind="office", read_path=side, text_path=side,
                      note=f"{ext} 문서에서 추출한 텍스트")
        elif ext in _TEXTY or not ext:
            entry(name=safe, kind="text", read_path=path, text_path=path)
        else:
            try:
                data.decode("utf-8")
                entry(name=safe, kind="text", read_path=path, text_path=path)
            except UnicodeDecodeError:
                entry(name=safe, kind="skip",
                      note="지원하지 않는 바이너리 형식(무시됨)")
    return workdir, manifest
