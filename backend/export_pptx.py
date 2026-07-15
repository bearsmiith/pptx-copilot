"""Export a Deck to .pptx via python-pptx.

Consumes the SAME DrawItems as the SVG renderer — preview and export match
by construction. Polygons (vias, underfill fillets) use freeform shapes.
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

from layout import (
    layout_slide, SLIDE_W, SLIDE_H,
    TextItem, FigureItem, RectShape, PolyShape, CircleShape, EdgeShape, Callout,
)
from models import Deck
from palette import material_colors, DARK_ROLES

NAVY = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x2F, 0x6F, 0xED)
GRAY = RGBColor(0x5B, 0x64, 0x72)
LEADER = RGBColor(0x8A, 0x93, 0xA0)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)


def _rgb(hexstr: str) -> RGBColor:
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _style_fill(shape, material: str, line_w: float = 1.5):
    fill_hex, stroke_hex = material_colors(material)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(stroke_hex)
    shape.line.width = Pt(line_w)


def _add_arrowhead(connector):
    try:
        ln = connector.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
        ln.append(tail)
    except Exception:
        pass


def _add_text(slide, t: TextItem):
    tb = slide.shapes.add_textbox(Inches(t.x), Inches(t.y), Inches(t.w), Inches(t.h))
    tf = tb.text_frame
    tf.word_wrap = True

    def style(p, role):
        if role == "title":
            p.font.size, p.font.bold, p.font.color.rgb = Pt(28), True, NAVY
        elif role == "subtitle":
            p.font.size, p.font.color.rgb = Pt(19), GRAY
        elif role == "caption":
            p.font.size, p.font.italic, p.font.color.rgb = Pt(12.5), True, GRAY
            p.alignment = PP_ALIGN.CENTER
        elif role == "panel_title":
            p.font.size, p.font.bold, p.font.color.rgb = Pt(15), True, NAVY
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.size, p.font.color.rgb = Pt(16), NAVY

    if t.role == "bullets":
        for i, b in enumerate(t.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.space_after = Pt(7)
            style(p, "body")
    else:
        p = tf.paragraphs[0]
        p.text = t.text
        style(p, t.role)


def _inner_label(shape, text: str, material: str, size_px: int):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(45720)   # 0.05"
    tf.margin_top = tf.margin_bottom = Emu(18288)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(max(9, int(size_px * 0.72)))   # px(96dpi) -> pt approx
    p.font.bold = True
    p.font.color.rgb = LIGHT if material in DARK_ROLES else NAVY


def _add_shape(slide, s):
    if isinstance(s, RectShape):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if s.rx > 0 else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(kind, Inches(s.x), Inches(s.y), Inches(s.w), Inches(s.h))
        if s.rx > 0:
            try:
                shp.adjustments[0] = min(0.5, s.rx / min(s.w, s.h))
            except Exception:
                pass
        _style_fill(shp, s.material)
        if s.inner_label:
            _inner_label(shp, s.inner_label, s.material, s.label_size)
    elif isinstance(s, CircleShape):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(s.cx - s.r), Inches(s.cy - s.r),
            Inches(2 * s.r), Inches(2 * s.r))
        _style_fill(shp, s.material, 1.2)
    elif isinstance(s, PolyShape):
        pts = s.points
        fb = slide.shapes.build_freeform(Emu(Inches(pts[0][0])), Emu(Inches(pts[0][1])))
        fb.add_line_segments([(Emu(Inches(px)), Emu(Inches(py))) for px, py in pts[1:]],
                             close=True)
        shp = fb.convert_to_shape()
        _style_fill(shp, s.material, 1.1)


def _add_callout(slide, c: Callout):
    # elbow leader as two connectors + small dot + label textbox
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(c.tx), Inches(c.ty),
        Inches(c.lx - 0.12), Inches(c.ly))
    conn.line.color.rgb = LEADER
    conn.line.width = Pt(1.1)
    conn2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(c.lx - 0.12), Inches(c.ly),
        Inches(c.lx), Inches(c.ly))
    conn2.line.color.rgb = LEADER
    conn2.line.width = Pt(1.1)
    dot_r = 0.03
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(c.tx - dot_r), Inches(c.ty - dot_r),
                                 Inches(2 * dot_r), Inches(2 * dot_r))
    dot.fill.solid()
    dot.fill.fore_color.rgb = LEADER
    dot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(c.lx), Inches(c.ly - 0.15), Inches(2.4), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True   # long labels wrap to a second line instead of clipping
    tf.margin_left = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = c.text
    p.font.size = Pt(11)
    p.font.color.rgb = NAVY


def _add_figure(slide, f: FigureItem):
    for e in f.edges:
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(e.x1), Inches(e.y1), Inches(e.x2), Inches(e.y2))
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(2.5)
        if getattr(e, "arrow", True):
            _add_arrowhead(conn)
    for s in f.shapes:
        _add_shape(slide, s)
    for c in f.callouts:
        _add_callout(slide, c)
    for t in f.texts:
        _add_text(slide, t)


def build_pptx(deck: Deck, overrides_by_slide=None) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for si, slide_plan in enumerate(deck.slides):
        slide = prs.slides.add_slide(blank)
        items = layout_slide(slide_plan)
        ov = (overrides_by_slide or {}).get(si) if isinstance(overrides_by_slide, dict) \
            else (overrides_by_slide[si] if overrides_by_slide
                  and si < len(overrides_by_slide) else None)
        if ov:
            from overrides import apply_overrides
            items = apply_overrides(items, ov)
        for it in items:
            if isinstance(it, TextItem):
                _add_text(slide, it)
            elif isinstance(it, FigureItem):
                _add_figure(slide, it)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
